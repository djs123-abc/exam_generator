"""
文档读取模块
支持: PDF, Word(.docx/.doc), Excel, PPT, 图片(OCR), TXT, URL网页
OCR 使用 PaddleOCR
"""
import os
import re
import logging
import urllib.request
import urllib.error
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    "pdf": [".pdf"],
    "word": [".docx", ".doc"],
    "excel": [".xlsx", ".xls", ".csv"],
    "ppt": [".pptx", ".ppt"],
    "image": [".jpg", ".jpeg", ".png", ".bmp", ".gif", ".tiff", ".tif", ".webp"],
    "text": [".txt", ".md", ".json"],
}

_ocr_instance = None


def get_ocr():
    global _ocr_instance
    if _ocr_instance is None:
        try:
            from paddleocr import PaddleOCR
            _ocr_instance = PaddleOCR(use_angle_cls=True, lang="ch", show_log=False)
            logger.info("PaddleOCR 初始化成功")
        except ImportError:
            logger.warning("PaddleOCR 未安装，图片OCR功能不可用")
    return _ocr_instance


def read_file(file_path: str) -> str:
    """读取单个文件，返回文本内容"""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    suffix = path.suffix.lower()

    for ftype, exts in SUPPORTED_EXTENSIONS.items():
        if suffix in exts:
            reader = {
                "pdf": read_pdf,
                "word": read_word,
                "excel": read_excel,
                "ppt": read_ppt,
                "image": read_image_ocr,
                "text": read_text,
            }.get(ftype)
            if reader:
                return reader(file_path)

    raise ValueError(f"不支持的文件格式: {suffix}")


def read_pdf(file_path: str) -> str:
    """读取 PDF 文件，自动 fallback 到 OCR"""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        texts = []
        ocr_pages = []
        for page_num, page in enumerate(doc):
            text = page.get_text("text").strip()
            if text and len(text) > 20:
                texts.append(f"[第{page_num+1}页]\n{text}")
            else:
                # 可能是图片型 PDF，走 OCR
                ocr_pages.append((page_num, page))

        if ocr_pages:
            ocr = get_ocr()
            if ocr:
                logger.info(f"PDF 有 {len(ocr_pages)} 页需要 OCR 识别")
                for page_num, page in ocr_pages:
                    # 渲染为图片
                    mat = fitz.Matrix(2.0, 2.0)  # 2x 分辨率
                    pix = page.get_pixmap(matrix=mat)
                    img_bytes = pix.tobytes("png")
                    result = ocr.ocr(img_bytes, cls=True)
                    if result and result[0]:
                        ocr_text = "\n".join(
                            [line[1][0] for line in result[0] if line and line[1]]
                        )
                        texts.append(f"[第{page_num+1}页(OCR)]\n{ocr_text}")
        doc.close()
        return "\n\n".join(texts)
    except Exception as e:
        logger.error(f"PDF 读取失败: {e}")
        raise


def read_word(file_path: str) -> str:
    """读取 Word 文档"""
    try:
        from docx import Document
        doc = Document(file_path)
        texts = []
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text)
        # 读取表格
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    texts.append(row_text)
        return "\n".join(texts)
    except Exception:
        # 尝试 docx2txt
        try:
            import docx2txt
            return docx2txt.process(file_path)
        except Exception as e:
            logger.error(f"Word 读取失败: {e}")
            raise


def read_excel(file_path: str) -> str:
    """读取 Excel/CSV 文件"""
    suffix = Path(file_path).suffix.lower()
    if suffix == ".csv":
        try:
            import chardet
            with open(file_path, "rb") as f:
                enc = chardet.detect(f.read())["encoding"] or "utf-8"
            with open(file_path, encoding=enc, errors="ignore") as f:
                return f.read()
        except Exception as e:
            logger.error(f"CSV 读取失败: {e}")
            raise
    else:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            texts = []
            for sheet in wb.worksheets:
                texts.append(f"[工作表: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) if c is not None else "" for c in row)
                    if row_text.strip(" |"):
                        texts.append(row_text)
            return "\n".join(texts)
        except Exception as e:
            logger.error(f"Excel 读取失败: {e}")
            raise


def read_ppt(file_path: str) -> str:
    """读取 PPT 文件"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [f"[幻灯片 {slide_num}]"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            texts.extend(slide_texts)
        return "\n".join(texts)
    except Exception as e:
        logger.error(f"PPT 读取失败: {e}")
        raise


def read_image_ocr(file_path: str) -> str:
    """图片 OCR 识别"""
    ocr = get_ocr()
    if not ocr:
        raise RuntimeError("PaddleOCR 未初始化，无法识别图片")
    try:
        result = ocr.ocr(file_path, cls=True)
        if not result or not result[0]:
            return ""
        lines = [line[1][0] for line in result[0] if line and line[1]]
        return "\n".join(lines)
    except Exception as e:
        logger.error(f"图片 OCR 失败: {e}")
        raise


def read_text(file_path: str) -> str:
    """读取纯文本文件"""
    try:
        import chardet
        with open(file_path, "rb") as f:
            raw = f.read()
        enc = chardet.detect(raw)["encoding"] or "utf-8"
        return raw.decode(enc, errors="ignore")
    except Exception as e:
        logger.error(f"文本文件读取失败: {e}")
        raise


def read_url(url: str) -> str:
    """抓取网页内容"""
    try:
        import urllib.request
        from bs4 import BeautifulSoup
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        }
        req = urllib.request.Request(url, headers=headers)
        with urllib.request.urlopen(req, timeout=30) as resp:
            raw = resp.read()
        # 检测编码
        import chardet
        enc = chardet.detect(raw)["encoding"] or "utf-8"
        html = raw.decode(enc, errors="ignore")
        soup = BeautifulSoup(html, "lxml")
        # 移除无用标签
        for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        # 压缩空行
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        return "\n".join(lines)
    except Exception as e:
        logger.error(f"URL 抓取失败: {e}")
        raise


def read_directory(dir_path: str, progress_callback=None) -> dict[str, str]:
    """
    递归读取目录下所有支持的文件
    返回 {文件名: 内容} 字典
    """
    results = {}
    all_exts = set()
    for exts in SUPPORTED_EXTENSIONS.values():
        all_exts.update(exts)

    files = []
    for root, dirs, filenames in os.walk(dir_path):
        # 忽略隐藏目录
        dirs[:] = [d for d in dirs if not d.startswith(".")]
        for fn in filenames:
            if Path(fn).suffix.lower() in all_exts:
                files.append(os.path.join(root, fn))

    for i, fp in enumerate(files):
        if progress_callback:
            progress_callback(i + 1, len(files), fp)
        try:
            content = read_file(fp)
            rel_path = os.path.relpath(fp, dir_path)
            results[rel_path] = content
            logger.info(f"读取成功: {rel_path} ({len(content)} 字符)")
        except Exception as e:
            logger.warning(f"读取失败 {fp}: {e}")
            results[os.path.relpath(fp, dir_path)] = f"[读取失败: {e}]"

    return results
